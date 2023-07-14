from fastapi import UploadFile, File
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import aspose.words as aw
import unicodedata
from models.question_answer import ask_questions
from interface.standard import QuestionsRequest


def clean_text(text: str):
    cleaned_text = []
    for char in text:
        if unicodedata.category(char)[0] != 'C':
            cleaned_text.append(char)
    return ''.join(cleaned_text)

def convert_pdf_to_text(url_file:Path)->str:
    pdf_reader = PyPDF2.PdfReader(url_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def convert_docx_to_text(url_file:Path)->str:
    document = Document(url_file)
    text = '\n'.join([paragraph.text for paragraph in document.paragraphs])
    return text

def convert_pptx_to_text(url_file:Path)->str:
    presentation = Presentation(url_file)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + "\n"
    return text

def convert_rtf_to_text(url_file:Path)->str:
    doc = aw.Document(str(url_file))
    text = doc.get_text()
    return text

def convert_text_to_object(text:str):
    questions = [
    "¿Cuál es la referencia?",
    "¿Qué tipo de documento es?",
    "¿Cuál es el número de expediente?",
    "¿Cuál es su tema principal?",
    "¿Cuál fue el magistrado?",
    "¿Qué sala es la responsable?",
    ]
    request = QuestionsRequest(text = text[:1000], questions= questions)
    return ask_questions(request)
     
